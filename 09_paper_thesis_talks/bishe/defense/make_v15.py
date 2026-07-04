# =============================================================
#  v13 -> v15: 浅色背景 / 黑字 / 字号≥16 / 去页脚 / 单列改版 / 参考文献页
# =============================================================
import re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

SRC = 'defense_slides_v13.pptx'
DST = 'defense_slides_v15.pptx'

# ----------------- Colors -----------------
BLACK    = RGBColor(0x00, 0x00, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xF4, 0xF7, 0xFB)   # 轻底
LIGHT2   = RGBColor(0xE8, 0xEE, 0xF6)   # 卡片底
TEAL     = RGBColor(0x1C, 0x72, 0x93)
AMBER    = RGBColor(0xF2, 0xA6, 0x5A)
GREY     = RGBColor(0xCF, 0xD8, 0xDC)

# ----------------- Filters -----------------
FOOTER_RE   = re.compile(r'杭州电子科技大学')
PAGENUM_RE  = re.compile(r'^\s*\d+\s*/\s*\d+\s*$')
DECOR_THIN  = 0.10  # inch；高度小于此值视为装饰条，保留原色

# ----------------- helpers -----------------
def luminance(rgb):
    try:
        return (0.299*rgb[0] + 0.587*rgb[1] + 0.114*rgb[2]) / 255.0
    except Exception:
        return 1.0

def is_dark(rgb):
    if rgb is None: return False
    return luminance(rgb) < 0.55

def set_slide_white(slide):
    """强制 slide-level 背景为白色（覆盖 dark 主题）。"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

def get_solid_rgb(shape):
    try:
        f = shape.fill
        if f.type == 1:  # solid
            return f.fore_color.rgb
    except Exception:
        return None
    return None

def lighten_fill(shape):
    """若 shape 当前为 solid 暗色填充，改为 LIGHT2。
       薄装饰条（h<0.1\"）保持原色（teal/amber accent 在白底依然好看）。
    """
    rgb = get_solid_rgb(shape)
    if rgb is None: return
    if not is_dark(rgb): return
    h_in = (shape.height or 0) / 914400
    w_in = (shape.width or 0)  / 914400
    if h_in < DECOR_THIN or w_in < DECOR_THIN:
        # 已是 dark thin bar -> 改为 teal 以与浅底对比
        shape.fill.fore_color.rgb = TEAL
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT2

def recolor_text_black(shape, min_pt=16):
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.color.rgb = BLACK
            except Exception:
                pass
            if run.font.size is None:
                continue
            if run.font.size.pt < min_pt:
                run.font.size = Pt(min_pt)

def strip_footer_shapes(slide):
    """删除 footer 文本 + 页码文本框。
       识别：文本含 '杭州电子科技大学' 或匹配 'N / M'。
    """
    to_del = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if FOOTER_RE.search(t) or PAGENUM_RE.fullmatch(t):
                to_del.append(sh)
    for sh in to_del:
        sh._element.getparent().remove(sh._element)

def delete_shapes(slide, shapes):
    for sh in shapes:
        sh._element.getparent().remove(sh._element)

# ----------------- Rebuild Slide 2 (Contents) — 单列 -----------------
OUTLINE_ITEMS = [
    ('01', '研究背景与意义',     '拓扑磁结构 · Hopfion · 信息载体'),
    ('02', '理论基础与构造方法',  'LLG 方程 · Hopf 指数 · 解析初始态'),
    ('03', '拓扑稳定性研究',      '阻挫交换体系 · 稳态尺寸 · 临界各向异性'),
    ('04', '自旋波驱动动力学',    '方向选择 · 频率响应 · 双向输运'),
    ('05', '神经形态计算应用',    'LIF 类比 · 双向可逆输运'),
]

def rebuild_slide_2(slide):
    """删除原 5 张卡片，重排为 5 行单列布局。"""
    # 保留：[0,1] 顶/底装饰条、[2] CONTENTS、[3] 答辩提纲、[4] 中间小条
    keep_names = {'Shape 0', 'Shape 1', 'Text 2', 'Text 3', 'Shape 4'}
    to_del = [sh for sh in slide.shapes if sh.name not in keep_names]
    delete_shapes(slide, to_del)

    # 将 “中间小条 (Shape 4)” 移到标题下方居左短横线
    for sh in slide.shapes:
        if sh.name == 'Shape 4':
            sh.left  = Inches(0.55)
            sh.top   = Inches(1.78)
            sh.width = Inches(0.6)
            sh.height = Emu(36000)  # 0.04"
            sh.fill.solid(); sh.fill.fore_color.rgb = TEAL
            try: sh.line.fill.background()
            except Exception: pass

    # 5 行
    y0   = 2.05
    rowH = 0.95   # 每行 0.95"
    gap  = 0.07
    full_w = 12.3
    x0   = 0.5

    for i, (num, title, sub) in enumerate(OUTLINE_ITEMS):
        y = y0 + i * (rowH + gap)
        # 卡片底（圆角）
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x0), Inches(y), Inches(full_w), Inches(rowH))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = GREY
        card.line.width = Pt(0.5)
        # 数字
        num_box = slide.shapes.add_textbox(Inches(x0+0.25), Inches(y+0.10),
                                           Inches(1.2), Inches(rowH-0.20))
        tf = num_box.text_frame; tf.margin_left=0; tf.margin_right=0
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = num
        r.font.name = 'Georgia'; r.font.size = Pt(44); r.font.bold = True
        r.font.color.rgb = TEAL
        # 标题
        t_box = slide.shapes.add_textbox(Inches(x0+1.7), Inches(y+0.10),
                                          Inches(full_w-2.0), Inches(0.45))
        tf = t_box.text_frame; tf.margin_left=0; tf.word_wrap=True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.name = '思源宋体 CN'
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = BLACK
        # 副标题
        s_box = slide.shapes.add_textbox(Inches(x0+1.7), Inches(y+0.55),
                                          Inches(full_w-2.0), Inches(0.35))
        tf = s_box.text_frame; tf.margin_left=0; tf.word_wrap=True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = sub
        r.font.name = '思源宋体 CN'
        r.font.size = Pt(16); r.font.color.rgb = BLACK


# ----------------- Rebuild Slide 22 (Conclusions) — 单列 -----------------
CONCLUSION_ITEMS = [
    ('01', '解析化拓扑构造工具链',
     '实现任意 Q_H 与反铁磁交替背景下 Hopfion 解析生成 + 三维可视化，覆盖 Q_H = 1, 2, 4。'),
    ('02', '稳定性参数判据',
     '阻挫交换体系定位 K_u1,c ∈ (52, 55) × 10³ J/m³，确认 R_eq = 2.60 nm 内秉稳态尺寸。'),
    ('03', '类脑功能映射',
     '凝练动力学特征与 LIF 神经元功能映射，提出自旋波驱动的三维类脑器件概念。'),
]

def rebuild_slide_22(slide):
    keep_names = {'Shape 0', 'Shape 1', 'Text 2', 'Text 3', 'Shape 4'}
    to_del = [sh for sh in slide.shapes if sh.name not in keep_names]
    delete_shapes(slide, to_del)

    for sh in slide.shapes:
        if sh.name == 'Shape 4':
            sh.left  = Inches(0.55)
            sh.top   = Inches(1.68)
            sh.width = Inches(0.6)
            sh.height = Emu(36000)
            sh.fill.solid(); sh.fill.fore_color.rgb = AMBER
            try: sh.line.fill.background()
            except Exception: pass

    y0   = 1.95
    rowH = 1.55
    gap  = 0.12
    full_w = 12.3
    x0   = 0.5

    for i, (num, title, body) in enumerate(CONCLUSION_ITEMS):
        y = y0 + i * (rowH + gap)
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x0), Inches(y), Inches(full_w), Inches(rowH))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = GREY; card.line.width = Pt(0.5)
        # 左侧 amber 竖条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x0), Inches(y),
                                      Inches(0.10), Inches(rowH))
        bar.fill.solid(); bar.fill.fore_color.rgb = AMBER
        try: bar.line.fill.background()
        except Exception: pass
        # 数字（大）
        num_box = slide.shapes.add_textbox(Inches(x0+0.25), Inches(y+0.15),
                                           Inches(1.4), Inches(rowH-0.30))
        tf = num_box.text_frame; tf.margin_left=0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = num
        r.font.name = 'Georgia'; r.font.size = Pt(56); r.font.bold=True
        r.font.color.rgb = TEAL
        # 标题
        t_box = slide.shapes.add_textbox(Inches(x0+1.9), Inches(y+0.20),
                                          Inches(full_w-2.2), Inches(0.55))
        tf = t_box.text_frame; tf.margin_left=0; tf.word_wrap=True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.name = '思源宋体 CN'
        r.font.size = Pt(24); r.font.bold=True; r.font.color.rgb = BLACK
        # body
        b_box = slide.shapes.add_textbox(Inches(x0+1.9), Inches(y+0.78),
                                          Inches(full_w-2.2), Inches(rowH-0.85))
        tf = b_box.text_frame; tf.margin_left=0; tf.word_wrap=True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = body
        r.font.name = '思源宋体 CN'
        r.font.size = Pt(17); r.font.color.rgb = BLACK


# ----------------- References Slide -----------------
REFERENCES = [
    'Sutcliffe P. Hopfions in chiral magnets [J]. J Phys A, 2018, 51: 375401.',
    'Wang X S, Qaiumzadeh A, Brataas A. Current-driven dynamics of magnetic hopfions [J]. Phys Rev Lett, 2019, 123: 147203.',
    'Zheng F, Kiselev N S, Rybakov F N, et al. Hopfion rings in a cubic chiral magnet [J]. Nature, 2023, 623: 718-723.',
    'Kent N, Reynolds N, Raftrey D, et al. Creation and observation of hopfions in magnetic multilayer systems [J]. Nat Commun, 2021, 12: 1562.',
    'Guslienko K. Magnetic hopfions: A review [J]. Magnetism, 2024, 4: 383-399.',
    'Knapman R, Tausendpfund T, Díaz S A, et al. Spacetime magnetic hopfions from skyrmion braiding [J]. Commun Phys, 2024, 7: 151.',
    'Göbel B, Mertig I, Tretiakov O A. Beyond skyrmions: Alternative magnetic quasiparticles [J]. Phys Rep, 2021, 895: 1-28.',
    'Vansteenkiste A, Leliaert J, Dvornik M, et al. The design and verification of MuMax3 [J]. AIP Adv, 2014, 4: 107133.',
]

def add_references_slide(pres):
    blank = pres.slide_layouts[0]
    s = pres.slides.add_slide(blank)
    set_slide_white(s)
    # 顶/底装饰条
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Emu(36000))
    top.fill.solid(); top.fill.fore_color.rgb = TEAL
    try: top.line.fill.background()
    except Exception: pass
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.46), Inches(13.333), Emu(36000))
    bot.fill.solid(); bot.fill.fore_color.rgb = TEAL
    try: bot.line.fill.background()
    except Exception: pass
    # 英文小标 + 中文大标
    en = s.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.45))
    p = en.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = 'REFERENCES'
    r.font.name='Georgia'; r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=TEAL

    cn = s.shapes.add_textbox(Inches(0.5), Inches(1.00), Inches(12.3), Inches(0.7))
    p = cn.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = '参考文献'
    r.font.name='思源宋体 CN'; r.font.size=Pt(36); r.font.bold=True; r.font.color.rgb=BLACK
    # 短分隔条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.80),
                              Inches(0.6), Emu(36000))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    try: bar.line.fill.background()
    except Exception: pass

    # 参考文献列表
    list_box = s.shapes.add_textbox(Inches(0.55), Inches(2.10),
                                     Inches(12.3), Inches(5.20))
    tf = list_box.text_frame; tf.word_wrap = True
    tf.margin_left = 0
    for i, ref in enumerate(REFERENCES):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(8)
        r = para.add_run(); r.text = f'[{i+1}] {ref}'
        r.font.name = '思源宋体 CN'
        r.font.size = Pt(16); r.font.color.rgb = BLACK
    return s


# ============================================================
def main():
    pres = Presentation(SRC)

    for idx, slide in enumerate(pres.slides):
        set_slide_white(slide)
        strip_footer_shapes(slide)
        # 全局：浅化暗填充 + 黑字 + ≥16pt
        for sh in list(slide.shapes):
            lighten_fill(sh)
            recolor_text_black(sh, min_pt=16)

    # 重建 slide 2 与 slide 22
    rebuild_slide_2(pres.slides[1])
    rebuild_slide_22(pres.slides[21])

    # 追加参考文献页（先 append 到尾部，再移到 22 与 23 之间）
    add_references_slide(pres)
    # 此时 References 是最后一张 (idx=23, 1-indexed 24)，Thanks 是 idx=22 (1-indexed 23)
    # 目标：References 成为 23, Thanks 成为 24
    xml = pres.slides._sldIdLst
    items = list(xml)
    ref_id = items[-1]
    xml.remove(ref_id)
    # 插入到 thanks 之前 → 原 idx=22 位置
    xml.insert(22, ref_id)

    pres.save(DST)
    print(f'Saved {DST}, total slides = {len(pres.slides)}')

if __name__ == '__main__':
    main()
