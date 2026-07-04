# =============================================================
#  v19 -> v20: 替换参考文献页 (8 条) 为论文完整 30 条 GB/T 7714 (分 3 页)
# =============================================================
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = 'defense_slides_v19.pptx'
DST = 'defense_slides_v20.pptx'
REFS_TXT = 'refs_final.txt'

# ---- Style ----
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL  = RGBColor(0x1C, 0x72, 0x93)
EA_FONT  = 'SimHei'
LAT_FONT = 'Calibri'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def set_run_font(run, ea=EA_FONT, lat=LAT_FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ('latin','ea','cs'):
        for el in rPr.findall(f'{{{NS_A}}}{tag}'):
            rPr.remove(el)
    for tag, face in [('latin', lat), ('ea', ea), ('cs', lat)]:
        el = etree.SubElement(rPr, f'{{{NS_A}}}{tag}')
        el.set('typeface', face)
    try: run.font.name = lat
    except Exception: pass

def load_refs():
    out = []
    with open(REFS_TXT, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line: continue
            # remove leading [N] (we'll re-add bold)
            m = re.match(r'^\[(\d+)\]\s*(.+)$', line)
            if m:
                out.append((int(m.group(1)), m.group(2)))
            else:
                out.append((None, line))
    # 修复 DOI 拆行造成的空格
    fixed = []
    for n, t in out:
        t = re.sub(r'10\.1007/BF014579\s*962', '10.1007/BF01457962', t)
        t = re.sub(r'PhysRevB\.104\.0\s*75102', 'PhysRevB.104.075102', t)
        t = re.sub(r'(\d)\s+(\d{3,})', lambda m: f'{m.group(1)}{m.group(2)}' if 'DOI' in t[max(0,t.find(m.group(0))-30):t.find(m.group(0))] else m.group(0), t)
        # collapse double space
        t = re.sub(r'\s+', ' ', t).strip()
        fixed.append((n, t))
    return fixed

def add_page_decor(s):
    """顶/底 teal 装饰条 — 与其它页保持视觉一致。"""
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Emu(36000))
    top.fill.solid(); top.fill.fore_color.rgb = TEAL
    try: top.line.fill.background()
    except Exception: pass
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.46), Inches(13.333), Emu(36000))
    bot.fill.solid(); bot.fill.fore_color.rgb = TEAL
    try: bot.line.fill.background()
    except Exception: pass

def add_title(s, page_idx, total_pages):
    """REFERENCES + 参考文献 (i/N) + 短青色分隔条"""
    # REFERENCES
    en = s.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.45))
    p = en.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = 'REFERENCES'
    set_run_font(r); r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=TEAL
    # 参考文献 (i/N)
    cn = s.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.7))
    p = cn.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = '参考文献'
    set_run_font(r); r.font.size=Pt(32); r.font.bold=True; r.font.color.rgb=BLACK
    r2 = p.add_run(); r2.text = f'  ({page_idx}/{total_pages})'
    set_run_font(r2); r2.font.size=Pt(20); r2.font.color.rgb=TEAL
    # 短青条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.78),
                              Inches(0.6), Emu(36000))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    try: bar.line.fill.background()
    except Exception: pass

def add_refs_to_slide(s, entries):
    """entries: list of (idx, text). 添加文献列表（一个 textbox 多段）。
       使用 hanging indent: 序号 [N] 与正文分离, 视觉对齐。
    """
    box = s.shapes.add_textbox(Inches(0.55), Inches(2.00),
                                Inches(12.30), Inches(5.40))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, (n, txt) in enumerate(entries):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(2)
        para.line_spacing = 1.0
        # paragraph properties: small hanging indent
        pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
        pPr.set('indent', '-457200')   # hanging
        pPr.set('marL',  '457200')     # 0.5 inch left margin
        # [N] in teal bold
        r1 = para.add_run(); r1.text = f'[{n}] '
        set_run_font(r1); r1.font.size=Pt(16); r1.font.bold=True; r1.font.color.rgb=TEAL
        # body
        r2 = para.add_run(); r2.text = txt
        set_run_font(r2); r2.font.size=Pt(16); r2.font.color.rgb=BLACK

def clear_slide_shapes(slide):
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)

def main():
    pres = Presentation(SRC)
    refs = load_refs()
    assert len(refs) == 30, f'expected 30 refs, got {len(refs)}'

    page_chunks = [refs[0:8], refs[8:15], refs[15:22], refs[22:30]]
    TOTAL_PAGES = len(page_chunks)

    # --- 1. 找到现有参考文献页（slide 26, idx=25）→ 清空 + 重写为第 1/N
    refs_slide = pres.slides[25]
    clear_slide_shapes(refs_slide)
    bg = refs_slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    add_page_decor(refs_slide)
    add_title(refs_slide, 1, TOTAL_PAGES)
    add_refs_to_slide(refs_slide, page_chunks[0])

    # --- 2. 末尾追加 N-1 张新页
    blank = pres.slide_layouts[0]
    for i, chunk in enumerate(page_chunks[1:], start=2):
        s = pres.slides.add_slide(blank)
        bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        add_page_decor(s)
        add_title(s, i, TOTAL_PAGES)
        add_refs_to_slide(s, chunk)

    # --- 3. 把刚追加的 N-1 张移动到 Thanks 之前
    xml = pres.slides._sldIdLst
    items = list(xml)
    n_added = TOTAL_PAGES - 1
    moved = items[-n_added:]
    for sid in moved:
        xml.remove(sid)
    # 定位 Thanks
    slides_iter = list(pres.slides)
    thanks_idx = None
    for k, slide in enumerate(slides_iter):
        for sh in slide.shapes:
            if sh.has_text_frame and 'THANKS' in sh.text_frame.text.upper():
                thanks_idx = k; break
        if thanks_idx is not None: break
    assert thanks_idx is not None, 'THANKS slide not found'
    for off, sid in enumerate(moved):
        xml.insert(thanks_idx + off, sid)

    pres.save(DST)
    print(f'Saved {DST}, total slides = {len(pres.slides)}')

if __name__ == '__main__':
    main()
