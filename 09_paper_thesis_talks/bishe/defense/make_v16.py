# =============================================================
#  v15 -> v16: 字体统一 + 删除正文页英文 phrase 副标题
# =============================================================
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree

SRC = 'defense_slides_v15.pptx'
DST = 'defense_slides_v16.pptx'

# 字体策略：中文 SimHei，西文 Calibri
EA_FONT  = 'SimHei'
LAT_FONT = 'Calibri'

# 需要整体删除的英文副标题（slide_1based -> 文本起始片段）
KILL_SUBTITLES = {
    4:  'From 2D Skyrmion',
    5:  'From Computing Bottleneck',
    6:  'State-of-the-Art Review',
    7:  'Research Roadmap',
    12: 'Drift Trajectory',
    13: 'Equilibrium Size',
    16: 'Polarization Selectivity',
    18: 'Bidirectional Transport',
    20: 'Mapping Hopfion Dynamics',
    21: 'Spin-Wave Excitation',
}

NSMAP_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def set_run_font(run):
    """同时设置 latin / ea / cs typeface，保证中英混排字体一致。"""
    rPr = run._r.get_or_add_rPr()
    # 移除旧的 typeface 节点
    for tag in ('latin', 'ea', 'cs'):
        for el in rPr.findall(f'{{{NSMAP_A}}}{tag}'):
            rPr.remove(el)
    # 写入新 typeface
    for tag, face in [('latin', LAT_FONT), ('ea', EA_FONT), ('cs', LAT_FONT)]:
        el = etree.SubElement(rPr, f'{{{NSMAP_A}}}{tag}')
        el.set('typeface', face)
    # 也通过 python-pptx API 同步 latin
    try:
        run.font.name = LAT_FONT
    except Exception:
        pass

def kill_subtitle(slide, prefix):
    """删除文本以 prefix 开头的 shape。"""
    to_del = []
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        if t.startswith(prefix):
            to_del.append(sh)
    for sh in to_del:
        sh._element.getparent().remove(sh._element)

def main():
    pres = Presentation(SRC)

    # 1) 删英文副标题
    for slide_no, prefix in KILL_SUBTITLES.items():
        kill_subtitle(pres.slides[slide_no - 1], prefix)

    # 2) 字体统一
    for s in pres.slides:
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    set_run_font(run)

    pres.save(DST)
    print(f'Saved {DST}')

if __name__ == '__main__':
    main()
