"""
v17 -> v18: 在 ch03 主页 (slide 9) 后插入 3 张细节页:
  9.1  理论: 局域旋转矩阵 + 环面坐标 -> 解析磁化场
  9.2  离散化: Python 程序流程 + AFM 棋盘格背景
  9.3  可视化: Marching Cubes + 面内相位色相 + AFM 自动解调

风格沿用 v17 (Calibri / Ocean Gradient / 白卡片+淡蓝描边).
"""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = '/mnt/d/Research/Hopfion/09_paper_thesis_talks/bishe/defense/defense_slides_v17.pptx'
DST = '/mnt/d/Research/Hopfion/09_paper_thesis_talks/bishe/defense/defense_slides_v19.pptx'

INSERT_AFTER = 9   # 在第 9 张 (Hopfion 解析构造主页) 之后插入

# Ocean Gradient
TEAL       = RGBColor(0x1C, 0x72, 0x93)
ICE_BLUE   = RGBColor(0xE8, 0xEE, 0xF6)
CARD_LINE  = RGBColor(0xCA, 0xDC, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
MUTED      = RGBColor(0x55, 0x66, 0x77)

# --- 标题区坐标 (与 v17 完全一致) ---
LBAR = (0.50, 0.55, 0.12, 0.70)
TITLE = (0.75, 0.45, 11.50, 0.55)
UNDER = (0.75, 1.10, 0.50, 0.04)


def add_rect(slide, x, y, w, h, fill, line_color=None, line_w_pt=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w_pt)
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, text, *, size=16, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    # 第一段
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, list):
        first = True
        for line in text:
            if first:
                para = p
                first = False
            else:
                para = tf.add_paragraph()
            para.alignment = align
            r = para.add_run(); r.text = line
            r.font.name = font; r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    else:
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_header(slide, title_text):
    # 左竖条
    bar = add_rect(slide, *LBAR, fill=ICE_BLUE, line_color=TEAL, line_w_pt=1.0)
    # 标题
    add_text(slide, *TITLE, text=title_text, size=30, bold=True,
             color=BLACK, align=PP_ALIGN.LEFT)
    # 下划线短条
    add_rect(slide, *UNDER, fill=TEAL, line_color=TEAL, line_w_pt=1.0)


def add_card(slide, x, y, w, h, num, head, body_lines):
    """白卡片: 左上角圆形序号 + 小标题 + 多行正文."""
    # 背景白卡
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=CARD_LINE, line_w_pt=0.5)
    # 序号圆 (teal 填色)
    cx, cy, cw, ch = x + 0.20, y + 0.20, 0.55, 0.55
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(cx), Inches(cy), Inches(cw), Inches(ch))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.name = 'Calibri'; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = WHITE
    # 小标题
    add_text(slide, x + 0.85, y + 0.22, w - 1.0, 0.45, head,
             size=17, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    # 正文（多行）
    add_text(slide, x + 0.25, y + 0.92, w - 0.45, h - 1.05, body_lines,
             size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT)


def add_bottom_note(slide, y, head, body):
    """底部宽卡片: 左侧装饰条 + 小标题 + 长正文."""
    add_rect(slide, 0.50, y, 12.33, 1.15, fill=WHITE,
             line_color=CARD_LINE, line_w_pt=0.5)
    add_rect(slide, 0.50, y, 0.18, 1.15, fill=ICE_BLUE,
             line_color=TEAL, line_w_pt=1.0)
    add_text(slide, 0.85, y + 0.05, 2.5, 0.40, head,
             size=16, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    add_text(slide, 0.85, y + 0.45, 11.80, 0.65, body,
             size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT)


def add_formula_block(slide, x, y, w, h, lines):
    """公式块: 白底淡边卡片 + 居中 Cambria Math 文字."""
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=CARD_LINE, line_w_pt=0.5)
    add_text(slide, x + 0.10, y + 0.10, w - 0.20, h - 0.20, lines,
             size=13, bold=False, color=TEAL, align=PP_ALIGN.CENTER,
             font='Cambria Math')


# =====================  Slide A: 理论构造  =====================
def build_slide_theory(p):
    slide = p.slides.add_slide(p.slide_layouts[0])  # only layout (blank)
    add_header(slide, '怎么"画"出一个 Hopfion？')

    cards = [
        (0.54, 1.40, 4.10, 4.05, 1, '起点：一片均匀磁体',
         ['先想象一整块磁体，',
          '所有"小磁针"都朝同一个方向。',
          '',
          '我们要在里面，',
          '"扭"出一个甜甜圈形状的纽结。']),
        (4.75, 1.40, 4.10, 4.05, 2, '关键：换个坐标',
         ['直接用 xyz 写公式太麻烦。',
          '',
          '换一套"绕甜甜圈"的坐标后，',
          '复杂的扭转变成两个简单的',
          '"转圈次数"。']),
        (8.96, 1.40, 4.10, 4.05, 3, '形状由两个整数决定',
         ['p = 沿管子绕几圈',
          'q = 围着甜甜圈绕几圈',
          '',
          '拓扑荷  Q_H = p × q。',
          '想换形状，只需改两个数字。']),
    ]
    for c in cards:
        add_card(slide, *c)

    add_bottom_note(slide, 5.65,
                    '一句话',
                    '把"甜甜圈纽结"的形状压缩成两个整数 (p, q)；改数字就能生成任意 Hopfion。')


# =====================  Slide B: 程序化离散 + AFM  =====================
def build_slide_discrete(p):
    slide = p.slides.add_slide(p.slide_layouts[0])
    add_header(slide, '怎么把公式"变成"仿真初始态？')

    cards = [
        (0.54, 1.40, 4.10, 4.05, 1, '三步走',
         ['① 把空间切成小方格',
          '② 每个格子算一下它的位置',
          '③ 套公式得到磁化方向',
          '',
          '一遍走完就拿到完整初始态。']),
        (4.75, 1.40, 4.10, 4.05, 2, '甜甜圈方向可调',
         ['默认甜甜圈平躺在地上。',
          '',
          '一个开关就能让它',
          '立起来、或斜着放。',
          '',
          '这样可以测不同方向',
          '自旋波的响应。']),
        (8.96, 1.40, 4.10, 4.05, 3, '也支持反铁磁',
         ['反铁磁 = 相邻格子方向相反。',
          '',
          '脚本里加一个开关，',
          '就能把"棋盘格"或"分层"',
          '的反铁磁背景叠到上面。']),
    ]
    for c in cards:
        add_card(slide, *c)

    add_bottom_note(slide, 5.65,
                    '输出',
                    '存成标准 OVF 文件，Mumax3 直接读取就能跑仿真。')


# =====================  Slide C: 可视化与解调  =====================
def build_slide_visualize(p):
    slide = p.slides.add_slide(p.slide_layouts[0])
    add_header(slide, '怎么把它"画"成 3D 图？')

    cards = [
        (0.54, 1.40, 4.10, 4.05, 1, '先画出甜甜圈外壳',
         ['在数据里挑出所有',
          '"小磁针刚好水平"的点。',
          '',
          '把这些点连起来，',
          '就得到一张曲面 —',
          '正好是甜甜圈的轮廓。']),
        (4.75, 1.40, 4.10, 4.05, 2, '用颜色显示方向',
         ['曲面上每个点的小磁针，',
          '方向不一样。',
          '',
          '红、黄、绿、蓝 转一圈，',
          '就能直观看到磁场',
          '怎么"绕"着甜甜圈转。']),
        (8.96, 1.40, 4.10, 4.05, 3, '反铁磁要先"抚平"',
         ['反铁磁里相邻格子方向相反，',
          '直接画曲面会"碎"成马赛克。',
          '',
          '先把翻转的格子翻回来，',
          '让数据看起来像普通铁磁，',
          '再画就正常了。']),
    ]
    for c in cards:
        add_card(slide, *c)

    add_bottom_note(slide, 5.65,
                    '验证',
                    '画出来的甜甜圈形状对、颜色平滑过渡 → 说明理论公式和程序都没错。')


# =====================  插入排序  =====================
def reorder(prs, after_idx, n_new):
    """python-pptx 没有官方 reorder, 直接改 sldIdLst."""
    sld_id_lst = prs.slides._sldIdLst
    items = list(sld_id_lst)
    # 末尾的 n_new 张是新加的; 把它们移到 after_idx 之后
    new_slides = items[-n_new:]
    head = items[:after_idx]   # 前 after_idx 张 (1..after_idx)
    tail = items[after_idx:-n_new]
    for el in items:
        sld_id_lst.remove(el)
    for el in head + new_slides + tail:
        sld_id_lst.append(el)


def main():
    prs = Presentation(SRC)
    build_slide_theory(prs)
    build_slide_discrete(prs)
    build_slide_visualize(prs)
    reorder(prs, after_idx=INSERT_AFTER, n_new=3)
    prs.save(DST)
    print(f'Saved {DST}')


if __name__ == '__main__':
    main()
