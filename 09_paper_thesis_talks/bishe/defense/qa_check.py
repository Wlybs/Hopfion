#!/usr/bin/env python3
"""程序化 PPT 审查：检测文本溢出、几何重叠、内容缺失。"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400

def emu_to_in(v):
    return v / EMU_PER_INCH

def text_dims_estimate(text, font_size_pt):
    """粗估文本所需宽高（英寸）。font_size_pt 单位 pt = 1/72 in。"""
    if not text:
        return 0, 0
    line_h_in = font_size_pt * 1.35 / 72.0          # 行高
    cn_chars = sum(1 for c in text if ord(c) > 127)
    en_chars = len(text) - cn_chars
    # 单字符宽度（英寸）
    cn_w = font_size_pt * 1.0 / 72.0
    en_w = font_size_pt * 0.55 / 72.0
    # 取最大行长度（按 \n 切分）
    max_line_w = 0
    for line in text.split("\n"):
        cn = sum(1 for c in line if ord(c) > 127)
        en = len(line) - cn
        max_line_w = max(max_line_w, cn*cn_w + en*en_w)
    return max_line_w, line_h_in

def shape_to_record(sp, slide_idx):
    """提取 shape 的位置、文字与字体信息。"""
    rec = {
        "slide": slide_idx,
        "name": sp.name,
        "type": str(sp.shape_type),
        "x": emu_to_in(sp.left or 0),
        "y": emu_to_in(sp.top or 0),
        "w": emu_to_in(sp.width or 0),
        "h": emu_to_in(sp.height or 0),
        "text": "",
        "fs": None,
        "issues": []
    }
    if sp.has_text_frame:
        tf = sp.text_frame
        rec["text"] = tf.text
        # 取首段首 run 的字号作为代表
        sizes = []
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(run.font.size.pt)
        rec["fs"] = max(sizes) if sizes else None
    return rec

def overlap(a, b, tol=0.02):
    """两个矩形是否相交（容差英寸）。"""
    return not (a["x"]+a["w"] <= b["x"]+tol or
                b["x"]+b["w"] <= a["x"]+tol or
                a["y"]+a["h"] <= b["y"]+tol or
                b["y"]+b["h"] <= a["y"]+tol)

def main():
    pptx_path = Path("/mnt/d/Research/Hopfion/09_paper_thesis_talks/bishe/defense/defense_slides.pptx")
    prs = Presentation(pptx_path)
    slide_w = emu_to_in(prs.slide_width)
    slide_h = emu_to_in(prs.slide_height)
    print(f"Slide size: {slide_w:.2f} x {slide_h:.2f} in")
    print(f"Slides: {len(prs.slides)}")
    print("=" * 80)

    total_issues = 0
    for i, slide in enumerate(prs.slides, 1):
        recs = []
        for sp in slide.shapes:
            try:
                recs.append(shape_to_record(sp, i))
            except Exception as e:
                pass
        # --- 检查 1: 元素是否超出幻灯片边界
        for r in recs:
            if r["x"] < -0.05 or r["y"] < -0.05:
                r["issues"].append(f"超出左上 ({r['x']:.2f},{r['y']:.2f})")
            if r["x"] + r["w"] > slide_w + 0.05:
                r["issues"].append(f"超出右边 ({r['x']+r['w']:.2f}>{slide_w:.2f})")
            if r["y"] + r["h"] > slide_h + 0.05:
                r["issues"].append(f"超出底部 ({r['y']+r['h']:.2f}>{slide_h:.2f})")

        # --- 检查 2: 文本溢出（估算）
        for r in recs:
            if r["text"] and r["fs"]:
                tw, th_per_line = text_dims_estimate(r["text"], r["fs"])
                lines = r["text"].count("\n") + 1
                # 简单估算：若单行宽 > 可用宽 * 1.05 且 高 < 1 行高，可能溢出
                if tw > r["w"] * 1.4 and r["h"] < th_per_line * 2:
                    r["issues"].append(f"文本可能横向溢出 (估{tw:.2f}>{r['w']:.2f})")
                # 多行垂直估算
                est_h = lines * th_per_line
                if est_h > r["h"] * 1.3 and r["h"] > 0.1:
                    r["issues"].append(f"文本可能纵向溢出 (估{est_h:.2f}>{r['h']:.2f})")

        # --- 检查 3: 文本框重叠（仅对有文本的）
        text_recs = [r for r in recs if r["text"].strip()]
        for a_idx in range(len(text_recs)):
            for b_idx in range(a_idx+1, len(text_recs)):
                a, b = text_recs[a_idx], text_recs[b_idx]
                if overlap(a, b, tol=0.03):
                    pair = f"重叠 ⟷ '{b['text'][:18]}…'@({b['x']:.2f},{b['y']:.2f})"
                    a["issues"].append(pair)

        # --- 报告
        slide_issues = [r for r in recs if r["issues"]]
        if slide_issues:
            print(f"\n[Slide {i}]  {len(slide_issues)} 项问题")
            for r in slide_issues:
                head = r["text"][:30].replace("\n"," | ") if r["text"] else f"[{r['type']}]"
                print(f"  · '{head}' @({r['x']:.2f},{r['y']:.2f},{r['w']:.2f}×{r['h']:.2f}) fs={r['fs']}")
                for iss in r["issues"]:
                    print(f"      ! {iss}")
                total_issues += 1
        else:
            print(f"[Slide {i}]  ✓ 无问题  ({len(recs)} 元素)")

    print("=" * 80)
    print(f"汇总: {total_issues} 项潜在问题")

if __name__ == "__main__":
    main()
