"""Build talk.pptx from talk_outline.md.

朴素布局：标题（无 P# 前缀）+ 正文文本 + 单张图（如有）+ 右下角小字"来源"。
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


HERE = Path(__file__).resolve().parent
OUTLINE = HERE / "talk_outline.md"
PPTX_OUT = HERE / "talk.pptx"


def split_pages(md: str):
    pages = []
    parts = re.split(r"\n## (P\d+ — [^\n]+)\n", md)
    for i in range(1, len(parts), 2):
        title = parts[i].strip()
        body = parts[i + 1].strip() if i + 1 < len(parts) else ""
        pages.append((title, body))
    return pages


def strip_p_prefix(title: str) -> str:
    return re.sub(r"^P\d+\s*—\s*", "", title).strip()


SKIP_PATTERNS = [
    re.compile(r"^\s*无图。?\s*$"),
    re.compile(r"^\s*无图（[^）]*）。?\s*$"),
    re.compile(r"^\s*无结论。?\s*$"),
    re.compile(r"^\s*仅描述现象，无结论。?\s*$"),
    re.compile(r"^\s*无图，无结论。?\s*$"),
]

PREFIX_STRIP = [
    re.compile(r"^无图。\s*"),
    re.compile(r"^无图（[^）]*）。\s*"),
    re.compile(r"^仅描述现象，无结论。\s*"),
    re.compile(r"^无图，无结论。\s*"),
]


def clean_markdown(text: str) -> str:
    lines = text.split("\n")
    out = []
    for line in lines:
        s = line.rstrip()
        if s.strip() in ("---", "___", "***"):
            continue
        s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
        s = re.sub(r"(?<!\*)\*([^*\n]+?)\*(?!\*)", r"\1", s)
        # Strip "无图" / "无结论" prefixes inline
        for pat in PREFIX_STRIP:
            s = pat.sub("", s)
        # Skip lines that are only "无图" / "无结论"
        if any(p.match(s) for p in SKIP_PATTERNS):
            continue
        out.append(s)
    # Collapse trailing empty lines
    while out and not out[-1].strip():
        out.pop()
    return "\n".join(out).strip()


def extract_image(body: str):
    m = re.search(r"!\[.*?\]\(([^)]+)\)", body)
    if not m:
        return None, body
    img_rel = m.group(1)
    img_path = HERE / img_rel if not img_rel.startswith("/") else Path(img_rel)
    body_no_img = re.sub(r"!\[.*?\]\([^)]+\)\s*\n*", "", body).strip()
    return img_path if img_path.exists() else None, body_no_img


def extract_source(body: str):
    """Pull out lines starting with 来源 (one or more)."""
    lines = body.split("\n")
    src_lines = []
    rest = []
    for line in lines:
        if line.lstrip().startswith("来源"):
            src_lines.append(line.strip())
        else:
            rest.append(line)
    return "\n".join(src_lines).strip(), "\n".join(rest).strip()


def add_slide(prs, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title (P# prefix stripped)
    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    p = tb_title.text_frame.paragraphs[0]
    p.text = strip_p_prefix(re.sub(r"\*\*(.+?)\*\*", r"\1", title))
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x2D, 0x5A)

    img_path, body_no_img = extract_image(body)
    body_clean = clean_markdown(body_no_img)
    source, body_final = extract_source(body_clean)

    # Body region
    if img_path:
        tb_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.4))
        slide.shapes.add_picture(str(img_path), Inches(8.3), Inches(1.4), width=Inches(4.7))
    else:
        tb_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.4))

    tf = tb_body.text_frame
    tf.word_wrap = True
    first = True
    for line in body_final.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line.rstrip()
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x20, 0x20, 0x20)

    # Source: bottom-right, small grey
    if source:
        tb_src = slide.shapes.add_textbox(Inches(5.0), Inches(6.95), Inches(8.0), Inches(0.4))
        tf = tb_src.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        # Collapse multi-line source into one or two short lines
        compact = " | ".join(s.strip() for s in source.split("\n") if s.strip())
        p.text = compact
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        p.font.italic = True


def main():
    md = OUTLINE.read_text(encoding="utf-8")
    pages = split_pages(md)
    print(f"Parsed {len(pages)} pages.")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for title, body in pages:
        add_slide(prs, title, body)

    prs.save(str(PPTX_OUT))
    print(f"Saved: {PPTX_OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
